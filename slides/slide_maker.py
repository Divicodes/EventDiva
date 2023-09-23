from pptx import Presentation
from pptx.util import Inches

# Create a Presentation object
prs = Presentation()

# Slide 1
slide_1 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_1.shapes.title
title.text = "Welcome to E-Diva – Turning Emails into Experiences!"
content_box = slide_1.placeholders[1]
content_box.text = "Brief Overview: E-Diva – Your Dynamic Event Management App\n"
content_box.text += "Problem: Overloaded with event emails? Missing out on experiences?"

# Slide 2
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_2.shapes.title
title.text = "E-Diva: Your Solution to Effortless Event Management!"
content_box = slide_2.placeholders[1]
content_box.text = "E-Diva Overview: Auto extraction of event info | Quick addition to the calendar\n"
content_box.text += "User Benefits: Time Efficiency, Never Miss Out, Seamless Experience\n"
content_box.text += "Security Reassurance: Committed to your privacy and data security"

# Slide 3
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_3.shapes.title
title.text = "Behind the Scenes – E-Diva’s Robust Architecture!"
content_box = slide_3.placeholders[1]
content_box.text = "Frontend: Sleek, User-Friendly Interface\n"
content_box.text += "Backend: Robust Services for Synchronization, Extraction, Notifications\n"
content_box.text += "Security: Encryption and Secure Authentication"

# Slide 4
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_4.shapes.title
title.text = "E-Diva in Action – A Live Demo!"
content_box = slide_4.placeholders[1]
content_box.text = "Experience the user-friendly interface and one-tap addition."

# Slide 5
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_5.shapes.title
title.text = "Ready to Transform Your Event Experience?"
content_box = slide_5.placeholders[1]
content_box.text = "Availability: E-Diva is Now Available on Android Mobile!\n"
content_box.text += "Final Statement: Transform Your Emails into Experiences with E-Diva!"

# Save the presentation
prs.save('EDiva_Presentation.pptx')
